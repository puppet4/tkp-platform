"""文档解析器模块。

支持多种文档格式的文本提取：
- PDF (.pdf)
- Word (.docx, .doc)
- PowerPoint (.pptx, .ppt)
- 纯文本 (.txt, .md, .json, .csv)
"""

import logging
from pathlib import Path
from typing import Protocol

logger = logging.getLogger("tkp_worker.parsers")


class DocumentParser(Protocol):
    """文档解析器协议。"""

    def can_parse(self, filename: str) -> bool:
        """判断是否支持解析该文件。"""
        ...

    def parse(self, file_bytes: bytes, filename: str) -> str:
        """解析文档并返回文本内容。"""
        ...


class PDFParser:
    """PDF 文档解析器。"""

    def can_parse(self, filename: str) -> bool:
        return filename.lower().endswith(".pdf")

    def parse(self, file_bytes: bytes, filename: str) -> str:
        try:
            from PyPDF2 import PdfReader
            from io import BytesIO
        except ImportError as exc:
            raise RuntimeError("PDF parsing requires 'pypdf2' package") from exc

        try:
            reader = PdfReader(BytesIO(file_bytes))
            text_parts = []
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{text}")

            result = "\n\n".join(text_parts)
            logger.info("parsed PDF: %s, pages=%d, chars=%d", filename, len(reader.pages), len(result))
            return result
        except Exception as exc:
            logger.exception("failed to parse PDF: %s", filename)
            raise RuntimeError(f"PDF parsing failed: {exc}") from exc


class WordParser:
    """Word 文档解析器 (.docx)。"""

    def can_parse(self, filename: str) -> bool:
        return filename.lower().endswith(".docx")

    def parse(self, file_bytes: bytes, filename: str) -> str:
        try:
            from docx import Document
            from io import BytesIO
        except ImportError as exc:
            raise RuntimeError("Word parsing requires 'python-docx' package") from exc

        try:
            doc = Document(BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            # 提取表格内容
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        tables_text.append(row_text)

            result = "\n\n".join(paragraphs)
            if tables_text:
                result += "\n\n--- Tables ---\n" + "\n".join(tables_text)

            logger.info("parsed Word: %s, paragraphs=%d, chars=%d", filename, len(paragraphs), len(result))
            return result
        except Exception as exc:
            logger.exception("failed to parse Word: %s", filename)
            raise RuntimeError(f"Word parsing failed: {exc}") from exc


class PowerPointParser:
    """PowerPoint 文档解析器 (.pptx)。"""

    def can_parse(self, filename: str) -> bool:
        return filename.lower().endswith(".pptx")

    def parse(self, file_bytes: bytes, filename: str) -> str:
        try:
            from pptx import Presentation
            from io import BytesIO
        except ImportError as exc:
            raise RuntimeError("PowerPoint parsing requires 'python-pptx' package") from exc

        try:
            prs = Presentation(BytesIO(file_bytes))
            slides_text = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_parts = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)

                if len(slide_parts) > 1:
                    slides_text.append("\n".join(slide_parts))

            result = "\n\n".join(slides_text)
            logger.info("parsed PowerPoint: %s, slides=%d, chars=%d", filename, len(prs.slides), len(result))
            return result
        except Exception as exc:
            logger.exception("failed to parse PowerPoint: %s", filename)
            raise RuntimeError(f"PowerPoint parsing failed: {exc}") from exc


class PlainTextParser:
    """纯文本解析器。"""

    SUPPORTED_EXTENSIONS = {".txt", ".md", ".json", ".csv", ".log", ".yaml", ".yml", ".xml", ".html"}

    def can_parse(self, filename: str) -> bool:
        return Path(filename).suffix.lower() in self.SUPPORTED_EXTENSIONS

    def parse(self, file_bytes: bytes, filename: str) -> str:
        try:
            # 尝试多种编码
            for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
                try:
                    text = file_bytes.decode(encoding)
                    logger.info("parsed text: %s, encoding=%s, chars=%d", filename, encoding, len(text))
                    return text
                except UnicodeDecodeError:
                    continue

            # 如果所有编码都失败，使用 utf-8 并忽略错误
            text = file_bytes.decode("utf-8", errors="ignore")
            logger.warning("parsed text with errors ignored: %s, chars=%d", filename, len(text))
            return text
        except Exception as exc:
            logger.exception("failed to parse text: %s", filename)
            raise RuntimeError(f"Text parsing failed: {exc}") from exc


class DocumentParserRegistry:
    """文档解析器注册表。"""

    def __init__(self):
        self.parsers: list[DocumentParser] = [
            PDFParser(),
            WordParser(),
            PowerPointParser(),
            PlainTextParser(),
        ]

    def parse(self, file_bytes: bytes, filename: str) -> str:
        """根据文件名选择合适的解析器并解析。"""
        for parser in self.parsers:
            if parser.can_parse(filename):
                return parser.parse(file_bytes, filename)

        raise ValueError(f"Unsupported file type: {filename}")

    def is_supported(self, filename: str) -> bool:
        """检查文件类型是否支持。"""
        return any(parser.can_parse(filename) for parser in self.parsers)


# 全局解析器实例
default_parser_registry = DocumentParserRegistry()